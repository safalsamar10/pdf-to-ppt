import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt

from image_prompt_generator import generate_image_prompt
from diagram_generator import create_flowchart


def build_ppt(slides):

    prs = Presentation()

    prs.slide_width = 12192000
    prs.slide_height = 6858000

    for slide_text in slides:

        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # TITLE
        slide.shapes.title.text = slide_text[:60]

        # CONTENT
        tf = slide.placeholders[1].text_frame
        tf.text = slide_text

        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.name = "Mangal"

        # AI IMAGE PROMPT BOX
        prompt = generate_image_prompt(slide_text)

        box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(5),
            Inches(8),
            Inches(1)
        )

        box.text_frame.text = "AI IMAGE: " + prompt

        # FLOWCHART IMAGE
        diagram = create_flowchart()

        slide.shapes.add_picture(
            diagram,
            Inches(6),
            Inches(1),
            Inches(3),
            Inches(3)
        )

    file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")

    prs.save(file.name)

    return file.name
