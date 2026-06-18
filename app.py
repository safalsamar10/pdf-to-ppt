import streamlit as st
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt
import tempfile

st.title("PDF → PPT Generator")

uploaded = st.file_uploader("Upload PDF", type=["pdf"])

def extract_text(pdf_file):
    doc = fitz.open(stream=pdf_file.read(), filetype="pdf")
    pages = []
    for page in doc:
        pages.append(page.get_text())
    return pages

def build_ppt(slides):
    prs = Presentation()
    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = s[:60]

        tf = slide.placeholders[1].text_frame
        tf.text = s

        p = tf.paragraphs[0]
        p.font.name = "Mangal"
        p.font.size = Pt(16)

    file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(file.name)
    return file.name

if uploaded:
    if st.button("Generate PPT"):
        pages = extract_text(uploaded)
        ppt_file = build_ppt(pages)

        with open(ppt_file, "rb") as f:
            st.download_button("Download PPT", f, file_name="output.pptx")
